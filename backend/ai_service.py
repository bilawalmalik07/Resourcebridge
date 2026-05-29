import os
import json
import requests
import tempfile
from pathlib import Path
from google import genai
from google.genai import types
from dotenv import load_dotenv

load_dotenv()

api_key = os.getenv("GEMINI_API_KEY")
client = genai.Client(api_key=api_key)

GEMINI_SUPPORTED = {
    ".pdf":  "application/pdf",
    ".png":  "image/png",
    ".jpg":  "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
    ".gif":  "image/gif",
    ".tiff": "image/tiff",
    ".tif":  "image/tiff",
    ".bmp":  "image/bmp",
}

TEXT_EXTRACTABLE = {
    ".docx", ".doc", ".xlsx", ".xls",
    ".pptx", ".ppt", ".txt", ".rtf", ".csv"
}

PROMPT = """
You are the core intelligence of ResourceBridge — a platform helping immigrant 
and working-class families understand important documents.

Analyze the provided document and perform FOUR tasks:

1. OCR_TEXT: Extract all readable text exactly as it appears.

2. SUMMARY_EN: Plain-language English summary covering:
   - What this document is
   - What action the family needs to take
   - Any deadlines
   - Who sent it and why it matters

3. SUMMARY_ES: Same summary in simple Spanish.

4. ACTION_ITEMS: JSON array of action items.
   Each: {"task": "...", "deadline": "YYYY-MM-DD or null", "priority": "high|medium|low"}
   If none: []

Return EXACTLY in this format:

---OCR_TEXT---
[text here]
---SUMMARY_EN---
[English summary here]
---SUMMARY_ES---
[Spanish summary here]
---ACTION_ITEMS---
[JSON array here]
"""


def _resolve_ext(file_url: str, content_type: str, original_filename: str | None) -> str | None:
    """Determine file extension from original filename, Content-Type, or URL (in priority order)."""
    if original_filename:
        ext = Path(original_filename).suffix.lower()
        if ext in GEMINI_SUPPORTED or ext in TEXT_EXTRACTABLE:
            print(f"Extension from original_filename: {ext}")
            return ext

    ct = content_type.lower()
    ct_map = {
        "pdf": ".pdf",
        "png": ".png",
        "jpeg": ".jpg", "jpg": ".jpg",
        "webp": ".webp",
        "gif": ".gif",
        "tiff": ".tiff",
        "bmp": ".bmp",
        "wordprocessingml": ".docx", "msword": ".docx",
        "spreadsheetml": ".xlsx", "ms-excel": ".xlsx",
        "presentationml": ".pptx", "powerpoint": ".pptx",
        "text/plain": ".txt",
        "text/csv": ".csv",
    }
    for keyword, ext in ct_map.items():
        if keyword in ct:
            print(f"Extension from Content-Type '{content_type}': {ext}")
            return ext

    url_ext = Path(file_url.split("?")[0]).suffix.lower()
    if url_ext in GEMINI_SUPPORTED or url_ext in TEXT_EXTRACTABLE:
        print(f"Extension from URL path: {url_ext}")
        return url_ext

    return None


def _extract_text_from_office(file_bytes: bytes, ext: str) -> str:
    if ext in (".docx", ".doc"):
        try:
            import mammoth
            import io
            return mammoth.extract_raw_text(io.BytesIO(file_bytes)).value
        except Exception:
            pass
        try:
            from docx import Document
            import io
            return "\n".join(p.text for p in Document(io.BytesIO(file_bytes)).paragraphs)
        except Exception:
            return ""

    elif ext in (".xlsx", ".xls"):
        try:
            import openpyxl
            import io
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    lines.append(
                        "\t".join(str(c) if c is not None else "" for c in row))
            return "\n".join(lines)
        except Exception:
            return ""

    elif ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(file_bytes))
            lines = []
            for i, slide in enumerate(prs.slides):
                lines.append(f"[Slide {i+1}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        lines.append(shape.text)
            return "\n".join(lines)
        except Exception:
            return ""

    elif ext in (".txt", ".csv", ".rtf"):
        return file_bytes.decode("utf-8", errors="ignore")

    return ""


def _send_file_to_gemini(file_bytes: bytes, ext: str) -> str:
    """Upload file to Gemini Files API and return response text."""
    mime_type = GEMINI_SUPPORTED[ext]
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        print(f"Uploading to Gemini Files API as {mime_type}...")
        uploaded = client.files.upload(
            file=tmp_path,
            config=types.UploadFileConfig(mime_type=mime_type)
        )
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[uploaded, PROMPT]
        )
        return response.text
    finally:
        os.unlink(tmp_path)


def _send_text_to_gemini(text: str) -> str:
    """Send extracted plain text to Gemini."""
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=f"{PROMPT}\n\nDocument text:\n\n{text}"
    )
    return response.text


def process_document_with_ai(file_url: str, original_filename: str | None = None) -> dict:
    try:
        print(f"Downloading: {file_url}")
        resp = requests.get(file_url, timeout=30)
        resp.raise_for_status()
        file_bytes = resp.content
        content_type = resp.headers.get("Content-Type", "")

        ext = _resolve_ext(file_url, content_type, original_filename)

        if not ext:
            print(f"Could not resolve extension — attempting UTF-8 plain text read")
            try:
                raw_text = file_bytes.decode("utf-8", errors="ignore").strip()
                if raw_text:
                    response_text = _send_text_to_gemini(raw_text)
                    return _parse_response(response_text)
            except Exception:
                pass
            return {
                "ocr_text": "",
                "ai_summary": "Could not determine file type. Please upload a PDF, image, Word, or Excel file.",
                "ai_summary_es": "No se pudo determinar el tipo de archivo. Suba un PDF, imagen, Word o Excel.",
                "action_items": [],
            }

        print(f"Processing as '{ext}'...")

        if ext in GEMINI_SUPPORTED:
            response_text = _send_file_to_gemini(file_bytes, ext)
        else:
            print(f"Extracting text from office format {ext}...")
            extracted = _extract_text_from_office(file_bytes, ext)
            if not extracted.strip():
                return {
                    "ocr_text": "",
                    "ai_summary": "Could not extract text from this file. Please convert to PDF and re-upload.",
                    "ai_summary_es": "No se pudo extraer texto. Por favor convierta a PDF.",
                    "action_items": [],
                }
            print(
                f"Extracted {len(extracted)} characters, sending to Gemini...")
            response_text = _send_text_to_gemini(extracted)

        return _parse_response(response_text)

    except requests.exceptions.RequestException as e:
        print(f"Download error: {e}")
        return {
            "ocr_text": "Failed to download the document.",
            "ai_summary": "Could not reach the file. Please try again.",
            "ai_summary_es": "No se pudo acceder al archivo.",
            "action_items": [],
        }
    except Exception as e:
        print(f"AI processing error: {type(e).__name__}: {e}")
        return {
            "ocr_text": f"Error: {type(e).__name__}: {str(e)[:300]}",
            "ai_summary": "Summary unavailable at this moment.",
            "ai_summary_es": "Resumen no disponible en este momento.",
            "action_items": [],
        }


def _parse_response(response_text: str) -> dict:
    ocr_text = _extract_section(response_text, "OCR_TEXT",    "SUMMARY_EN")
    summary_en = _extract_section(response_text, "SUMMARY_EN",  "SUMMARY_ES")
    summary_es = _extract_section(response_text, "SUMMARY_ES",  "ACTION_ITEMS")
    action_raw = _extract_section(response_text, "ACTION_ITEMS", None)
    try:
        action_items = json.loads(action_raw)
    except Exception:
        action_items = []
    return {
        "ocr_text":      ocr_text or "Text extraction processed.",
        "ai_summary":    summary_en or "Summary unavailable.",
        "ai_summary_es": summary_es or "Resumen no disponible.",
        "action_items":  action_items,
    }


def _extract_section(text: str, start_tag: str, end_tag: str | None) -> str:
    start_marker = f"---{start_tag}---"
    try:
        start_idx = text.index(start_marker) + len(start_marker)
        if end_tag:
            end_idx = text.index(f"---{end_tag}---")
            return text[start_idx:end_idx].strip()
        return text[start_idx:].strip()
    except ValueError:
        return ""
